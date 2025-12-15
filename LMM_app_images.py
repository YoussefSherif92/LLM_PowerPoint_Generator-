import google.generativeai as genai
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ==================================================
# 1. CONFIGURE GEMINI (OPTION B)
# ==================================================
genai.configure(api_key="AIzaSyA71rkTGD7XGCvyBSfO5Hm3BnIBc0pQ6Ts")
model = genai.GenerativeModel("models/gemini-flash-latest")

# ==================================================
# 2. PROMPT
# ==================================================
topic = "AI startup for healthcare"

prompt = f"""
Create a professional startup pitch PowerPoint about: {topic}

Return ONLY valid JSON in this format:
{{
  "title": "Presentation title",
  "slides": [
    {{
      "title": "Short professional slide title (no numbering)",
      "content": [
        "concise bullet point",
        "concise bullet point",
        "concise bullet point"
      ]
    }}
  ]
}}

Make exactly 5 slides.
Output JSON only.
"""

# ==================================================
# 3. CALL LLM (ROBUST JSON EXTRACTION)
# ==================================================
response = model.generate_content(prompt)
raw = response.text.strip()

start = raw.find("{")
end = raw.rfind("}")
if start == -1 or end == -1:
    raise RuntimeError("Gemini did not return valid JSON")

data = json.loads(raw[start:end + 1])

# ==================================================
# 4. LOAD IMAGES
# ==================================================
IMAGE_FOLDER = "images_powerpoint"

image_files = sorted([
    os.path.join(IMAGE_FOLDER, f)
    for f in os.listdir(IMAGE_FOLDER)
    if f.lower().endswith((".png", ".jpg", ".jpeg"))
])

# ==================================================
# 5. CREATE PRESENTATION
# ==================================================
prs = Presentation()

# ---------- GLOBAL STYLE ----------
BG_COLOR = RGBColor(245, 247, 250)
TITLE_COLOR = RGBColor(30, 64, 175)
TEXT_COLOR = RGBColor(55, 65, 81)

# ==================================================
# 6. TITLE SLIDE (CUSTOM – NO PLACEHOLDERS)
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_COLOR

title_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(2.5), Inches(7), Inches(1.5)
)
tf = title_box.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = data["title"]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.alignment = PP_ALIGN.CENTER

# ==================================================
# 7. CONTENT SLIDES (HUMAN DESIGN)
# ==================================================
for i, s in enumerate(data["slides"]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR

    # ---------- Title ----------
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(8.5), Inches(1)
    )
    tf = title_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = s["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # ---------- Layout Logic ----------
    image_on_right = (i % 2 == 0)

    text_left = Inches(0.75) if image_on_right else Inches(5.25)
    image_left = Inches(5.25) if image_on_right else Inches(0.75)

    # ---------- Text ----------
    text_box = slide.shapes.add_textbox(
        text_left, Inches(1.7), Inches(4), Inches(3.5)
    )
    tf = text_box.text_frame
    tf.clear()

    for bullet in s["content"]:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0

    # ---------- Image ----------
    if i < len(image_files):
        slide.shapes.add_picture(
            image_files[i],
            image_left,
            Inches(1.7),
            width=Inches(4)
        )

# ==================================================
# 8. SAVE
# ==================================================
prs.save("generated_presentation_human_style.pptx")
print("✅ Human-style PowerPoint created successfully.")
