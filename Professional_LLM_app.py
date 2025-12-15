import google.generativeai as genai
import requests
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ================================
# CONFIG
# ================================
GENAI_API_KEY = " AIzaSyA71rkTGD7XGCvyBSfO5Hm3BnIBc0pQ6Ts"
UNSPLASH_ACCESS_KEY = "-MwoMipMvDKhOPX9iIKPagzn7IOxPzywI39be-rjivg"

genai.configure(api_key=GENAI_API_KEY)
model = genai.GenerativeModel("models/gemini-flash-latest")

# ================================
# PROMPT
# ================================
prompt = """
You are a professional pitch-deck designer.

Create a PowerPoint about: AI startup for healthcare

Return ONLY valid JSON:
{
  "title": "Deck title",
  "theme": {
    "background": "light or dark",
    "primary_color": "#HEX",
    "accent_color": "#HEX",
    "font": "PowerPoint font"
  },
  "slides": [
    {
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2"],
      "image_query": "keywords",
      "layout": "text_left_image_right"
    }
  ]
}

Rules:
- Exactly 5 slides
- No markdown
- No explanation
"""

# ================================
# SAFE JSON EXTRACTION
# ================================
raw = model.generate_content(prompt).text.strip()
match = re.search(r"\{.*\}", raw, re.DOTALL)
data = json.loads(match.group())

theme = data["theme"]

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

PRIMARY = hex_to_rgb(theme["primary_color"])
ACCENT = hex_to_rgb(theme["accent_color"])
FONT_NAME = theme["font"]

# ================================
# UNSPLASH IMAGE FETCH
# ================================
def fetch_image(query, idx):
    url = "https://api.unsplash.com/search/photos"
    headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
    params = {"query": query, "per_page": 1}

    r = requests.get(url, headers=headers, params=params)
    img_url = r.json()["results"][0]["urls"]["regular"]

    path = f"img_{idx}.jpg"
    with open(path, "wb") as f:
        f.write(requests.get(img_url).content)
    return path

# ================================
# CREATE PRESENTATION
# ================================
prs = Presentation()

# -------- TITLE SLIDE --------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = data["title"]

tf = slide.shapes.title.text_frame
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.name = FONT_NAME
tf.paragraphs[0].font.color.rgb = PRIMARY

# -------- CONTENT SLIDES --------
for i, s in enumerate(data["slides"]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # ---- Title ----
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = s["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.name = FONT_NAME
    p.font.color.rgb = PRIMARY
    p.font.bold = True

    layout = s["layout"]

    def add_text(left):
        box = slide.shapes.add_textbox(
            left, Inches(1.4), Inches(4.5), Inches(4)
        )
        tf = box.text_frame
        tf.clear()
        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
            p.font.name = FONT_NAME

    def add_image(left):
        img = fetch_image(s["image_query"], i)
        slide.shapes.add_picture(
            img, left, Inches(1.4), width=Inches(4)
        )

    if layout == "text_left_image_right":
        add_text(Inches(0.5))
        add_image(Inches(5.3))
    elif layout == "image_left_text_right":
        add_image(Inches(0.5))
        add_text(Inches(5.3))
    elif layout == "full_image_with_caption":
        img = fetch_image(s["image_query"], i)
        slide.shapes.add_picture(img, Inches(0), Inches(1.2), width=Inches(10))
    elif layout == "text_only":
        add_text(Inches(2.5))

# ================================
# SAVE
# ================================
prs.save("AI_Themed_Professional_Presentation.pptx")
print("✅ THEMED, PERFECTLY ALIGNED PRESENTATION CREATED")
