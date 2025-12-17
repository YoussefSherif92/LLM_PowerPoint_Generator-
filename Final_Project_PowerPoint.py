import streamlit as st
import google.generativeai as genai
import requests
import json
import re
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ================================
# API KEYS (DIRECT)
# ================================
GENAI_API_KEY = "AIzaSyCKpaMKoDDm57MQv-W90uuLz7nR6AJeHag"
UNSPLASH_ACCESS_KEY = "-MwoMipMvDKhOPX9iIKPagzn7IOxPzywI39be-rjivg"

genai.configure(api_key=GENAI_API_KEY)
model = genai.GenerativeModel("models/gemini-flash-latest")

# ================================
# HELPER FUNCTIONS
# ================================
def extract_json(text):
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if not match:
        raise ValueError("No JSON returned")
    return json.loads(match.group())

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def fetch_image(query, idx):
    url = "https://api.unsplash.com/search/photos"
    headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
    params = {"query": query, "per_page": 1}
    r = requests.get(url, headers=headers, params=params)
    img_url = r.json()["results"][0]["urls"]["regular"]
    path = f"img_{idx}.jpg"
    with open(path, "wb") as f:
        f.write(requests.get(img_url).content)
    return path

# ================================
# POWERPOINT GENERATOR
# ================================
def generate_presentation(topic):
    prompt = f"""
You are a professional pitch-deck designer.

Create a 5-slide PowerPoint about: {topic}

Return ONLY JSON:
{{
  "title": "Deck title",
  "theme": {{
    "primary_color": "#HEX",
    "accent_color": "#HEX",
    "font": "Calibri",
    "design_style": "card | split | minimal"
  }},
  "slides": [
    {{
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "image_query": "professional photo keywords",
      "layout": "text_left_image_right | image_left_text_right | text_only"
    }}
  ]
}}

Rules:
- Exactly 5 slides
- No markdown
- No explanation
"""

    raw = model.generate_content(prompt).text
    data = extract_json(raw)

    PRIMARY = hex_to_rgb(data["theme"]["primary_color"])
    ACCENT = hex_to_rgb(data["theme"]["accent_color"])
    FONT = data["theme"]["font"]
    STYLE = data["theme"]["design_style"]

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(7), Inches(1.5)
    )
    tf = title.text_frame
    tf.text = data["title"]
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.name = FONT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY

    # Content slides
    for i, s in enumerate(data["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Design background
        if STYLE == "card":
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.4), Inches(1.1),
                Inches(9.5), Inches(5)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = ACCENT
            bg.line.fill.background()

        # Title
        title = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(9), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = s["title"]
        tf.paragraphs[0].font.size = Pt(26)
        tf.paragraphs[0].font.name = FONT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = PRIMARY

        layout = s["layout"]

        def add_text(left):
            box = slide.shapes.add_textbox(
                left, Inches(1.4), Inches(4.5), Inches(4)
            )
            tf = box.text_frame
            tf.clear()
            for b in s["bullets"]:
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
                p.font.name = FONT

        def add_image(left):
            img = fetch_image(s["image_query"], i)
            slide.shapes.add_picture(
                img, left, Inches(1.4), width=Inches(4)
            )

        if layout == "image_left_text_right":
            add_image(Inches(0.5))
            add_text(Inches(5.3))
        elif layout == "text_only":
            add_text(Inches(2.5))
        else:
            add_text(Inches(0.5))
            add_image(Inches(5.3))

    filename = f"AI_Presentation_{int(time.time())}.pptx"
    prs.save(filename)
    return filename

# ================================
# STREAMLIT UI
# ================================
st.set_page_config(page_title="AI PowerPoint Generator", layout="centered")

st.title("🧠 AI PowerPoint Generator")
st.write("Enter a topic and generate a **professional presentation**.")

topic = st.text_input("Presentation topic", placeholder="e.g. AI startup for healthcare")

if st.button("Generate Presentation"):
    if not topic:
        st.warning("Please enter a topic.")
    else:
        with st.spinner("Generating presentation..."):
            pptx_file = generate_presentation(topic)

        with open(pptx_file, "rb") as f:
            st.download_button(
                label="⬇️ Download PowerPoint",
                data=f,
                file_name=pptx_file,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        st.success("Presentation ready!")
