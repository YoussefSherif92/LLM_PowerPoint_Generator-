import google.generativeai as genai
import json
from pptx import Presentation

# -----------------------------
# Configure Gemini
# -----------------------------
genai.configure(api_key="AIzaSyA71rkTGD7XGCvyBSfO5Hm3BnIBc0pQ6Ts")

# ✅ CORRECT MODEL
model = genai.GenerativeModel("models/gemini-flash-latest")

# -----------------------------
# Prompt
# -----------------------------
topic = "AI startup for healthcare"

prompt = f"""
Create a PowerPoint about: {topic}

Return ONLY valid JSON in this format:
{{
  "title": "Presentation title",
  "slides": [
    {{
      "title": "Slide title",
      "content": ["bullet 1", "bullet 2"]
    }}
  ]
}}

Make exactly 5 slides.
Do not explain.
Do not use markdown.
Output JSON only.
"""

# -----------------------------
# Call Gemini
# -----------------------------
response = model.generate_content(prompt)

text = response.text
data = json.loads(text)

# -----------------------------
# Create PowerPoint
# -----------------------------
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = data["title"]

for s in data["slides"]:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = s["title"]
    slide.placeholders[1].text = "\n".join(s["content"])

prs.save("generated_presentation.pptx")
print("Presentation created!")
