import google.generativeai as genai
import requests
import json
import re
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ================================
# 1. API KEYS (DIRECT — AS YOU WANT)
# ================================
GENAI_API_KEY = "AIzaSyA71rkTGD7XGCvyBSfO5Hm3BnIBc0pQ6Ts"
UNSPLASH_ACCESS_KEY = "-MwoMipMvDKhOPX9iIKPagzn7IOxPzywI39be-rjivg"

genai.configure(api_key=GENAI_API_KEY)
model = genai.GenerativeModel("models/gemini-flash-latest")


# ================================
# 2. PROMPT (AI CONTROLS DESIGN)
# ================================
topic = "AI startup for healthcare"

prompt = f"""
You are a professional pitch-deck designer.

Create a 5-slide PowerPoint about: {topic}

Return ONLY valid JSON:

{{
  "title": "Deck title",
  "theme": {{
    "background": "light or dark",
    "primary_color": "#HEX",
    "accent_color": "#HEX",
    "font": "Calibri",
    "design_style": "card | split | minimal"
  }},
  "slides": [
    {{
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "image_query": "professional photo keywords",
      "layout": "text_left_image_right | image_left_text_right | full_image_with_caption | text_only",
      "caption": "optional"
    }}
  ]
}}

Rules:
- Exactly 5 slides
- No markdown
- No explanation
"""

# ================================
# 3. SAFE JSON EXTRACTION
# ================================
raw = model.generate_content(prompt).text.strip()
match = re.search(r"\{.*\}", raw, re.DOTALL)

if not match:
    print("❌ Gemini did not return JSON")
    print(raw)
    exit()

data = json.loads(match.group())


# ================================
# 4. THEME HELPERS
# ================================
def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

theme = data["theme"]
PRIMARY = hex_to_rgb(theme["primary_color"])
ACCENT = hex_to_rgb(theme["accent_color"])
FONT = theme["font"]
STYLE = theme["design_style"]


# ================================
# 5. UNSPLASH IMAGE SEARCH
# ================================
def fetch_image(query, idx):
    url = "https://api.unsplash.com/search/photos"
    headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
    params = {"query": query, "per_page": 1}

    r = requests.get(url, headers=headers, params=params)
    img_url = r.json()["results"][0]["urls"]["regular"]

    path = f"image_{idx}.jpg"
    with open(path, "wb") as f:
        f.write(requests.get(img_url).content)

    return path


# ================================
# 6. DESIGN BACKGROUNDS
# ================================
def draw_card(slide):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.2),
        Inches(9.5), Inches(5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT
    card.line.fill.background()


def draw_split(slide):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(3.8), Inches(7.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = ACCENT
    panel.line.fill.background()


def draw_minimal(slide):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(1.5), Inches(0.12)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()


# ================================
# 7. CREATE PRESENTATION
# ================================
prs = Presentation()

# ---- Title slide ----
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(2.5), Inches(7), Inches(1.5)
)
tf = title_box.text_frame
tf.text = data["title"]
tf.paragraphs[0].font.size = Pt(42)
tf.paragraphs[0].font.name = FONT
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PRIMARY


# ---- Content slides ----
for i, s in enumerate(data["slides"]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    if STYLE == "card":
        draw_card(slide)
    elif STYLE == "split":
        draw_split(slide)
    else:
        draw_minimal(slide)

    # Title
    title = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(9), Inches(0.8)
    )
    tf = title.text_frame
    tf.text = s["title"]
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.name = FONT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY

    layout = s["layout"]

    def add_text(left):
        box = slide.shapes.add_textbox(
            left, Inches(1.4), Inches(4.5), Inches(4)
        )
        tf = box.text_frame
        tf.clear()
        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
            p.font.name = FONT

    def add_image(left):
        img = fetch_image(s["image_query"], i)
        slide.shapes.add_picture(
            img, left, Inches(1.4), width=Inches(4)
        )

    if layout == "image_left_text_right":
        add_image(Inches(0.5))
        add_text(Inches(5.3))
    elif layout == "full_image_with_caption":
        img = fetch_image(s["image_query"], i)
        slide.shapes.add_picture(
            img, Inches(0.5), Inches(1.4), width=Inches(9)
        )
    elif layout == "text_only":
        add_text(Inches(2.5))
    else:
        add_text(Inches(0.5))
        add_image(Inches(5.3))


# ================================
# 8. SAVE (NO PERMISSION ERROR)
# ================================
filename = f"AI_Designed_Presentation_{int(time.time())}.pptx"
prs.save(filename)
print("✅ Presentation created:", filename)
